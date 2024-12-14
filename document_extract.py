"""This script extract text from various format document to make llm fine tunning training data
"""
import re
import os
import pathlib
import yaml
from minio import Minio
import pandas as pd
import pdfplumber
from langchain.document_loaders import TextLoader
from pptx import Presentation
from docx import Document
from docx.document import Document as _Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import _Cell, Table, _Row
from docx.text.paragraph import Paragraph
from HwpParser import HWPExtractor


class TextExtract:
    """
    Document extration main class
    """
    def __init__(self, bucket_name):
        with open('config/set-dev.yaml', "r", encoding="utf-8") as f:
            config = yaml.load(f, Loader=yaml.FullLoader)
            minio_info = config['minio']
        self._minio_address =  minio_info['address']
        self._accesskey = minio_info['accesskey']
        self._secretkey = minio_info['secretkey']
        self._minio_client = Minio(self._minio_address,
                                   access_key=self._accesskey,
                                   secret_key=self._secretkey,
                                   secure=False)
        self._bucket_name = bucket_name
        self.del_table = True

    def is_table_of_contents(self, page_text):
        """
        특정 페이지가 목차인지 판단
        :param page_text: 페이지에서 추출한 텍스트
        :return: 목차 페이지 여부 (True/False)
        """
        # 목차 페이지의 일반적인 패턴들
        patterns = [
            r'\.{5,}\s*\d+',  # 예: "Chapter 1 ...... 12"
            r'\·{5,}\d+',
            r'\b\d+\s*\.\s*\d+',  # 예: "1.1 Introduction"
            r'(Table of Contents|목차)',  # "Table of Contents" 또는 "목차"
            r'표 목 차',  # "Table of Contents" 또는 "목차"
        ]
        for pattern in patterns:
            if re.search(pattern, page_text, re.IGNORECASE):
                return True
        return False

    def get_context_pdffile_by_plumber(self, source_file_name:str)-> list:
        """
        get context from pdf file
        :param source_file_name: source pdf file path
        :return: extracted context text
        """
        print(f"Source path {source_file_name}")
        try:
            pdfplumb = pdfplumber.open(source_file_name)
            whole_page_extractinfo = ""
            page_exist_tbl = False
        except IOError as e:
            print(f"I/O error({e.errno}): {e.strerror}")
            return None

        total_pages = len(pdfplumb.pages)

        for page_num in range(1, total_pages):
            page_plumb_contents = {}
            table_list = []

            # pil_img = pdfplumb.pages[page_num].to_image(resolution=1200)
            # if page_num == 12:
            #     pil_img.save(f'{source_file_name}-{page_num}.png',"PNG", quantize=False)
            whole_text_a_page = ""
            if self.is_table_of_contents(pdfplumb.pages[page_num].extract_text()) is True:
                continue

            for table_info in pdfplumb.pages[page_num].find_tables():
                x0 = table_info.bbox[0]
                y0 = table_info.bbox[1]
                x1 = table_info.bbox[2]
                y1 = table_info.bbox[3]
                table_list.append((x0, y0, x1, y1))
                table = table_info.extract()
                df = pd.DataFrame(table[1::], columns=table[0])
                df.replace('\x00', '', inplace=True)
                df.replace('Ÿ', '*', inplace=True)
                page_plumb_contents[int(y0)] = {"type":"table",
                                                "value": df.to_markdown()}
            for content in pdfplumb.pages[page_num].extract_text_lines():
                txt_content = content['text']
                x0 = content['x0']
                y0 = content['top']
                x1 = content['x1']
                y1 = content['bottom']
                try:
                    if len(table_list) > 0:
                        if (table_list[0][0] < x0 and table_list[0][1] < y0 and
                                table_list[0][2] > x1 and table_list[0][3] > y1):
                            """
                            Filter context in outbound detected table contents
                            """
                            pass
                        else:
                            page_plumb_contents[int(y0)] = {"type": "text", "value":  txt_content}
                    else:
                        if y0 < 64.857: continue
                        if y0 >= 764.860: #Footer skeep
                            print(f"Page Footer Skeep : {x0} {y0} {txt_content}")
                            continue
                        if x0 > 200: #Innter title
                            print(f"Inner title  {x0} {y0} {txt_content}")
                            continue
                        page_plumb_contents[int(y0)] = {"type": "text", "value": txt_content}

                except Exception as e:
                    print(str(e))

            if len(page_plumb_contents) > 0:
                #각 페이지 단위 콘텐츠 결합
                pos_list = list(page_plumb_contents.keys())
                pos_list = sorted(pos_list)
                page_exist_tbl = False
                page_textonly_filtering = ""

                for position in pos_list:
                    line_text = page_plumb_contents[position]["value"]
                    if page_plumb_contents[position]["type"] == "table":
                        if self.del_table == False:
                            page_textonly_filtering = re.sub(r"(?<![\.\?\!])\n", " ", page_textonly_filtering)
                            whole_text_a_page += page_textonly_filtering + "\n" + line_text + "\n"
                        page_textonly_filtering = ""
                        page_exist_tbl = True
                    else:
                        if line_text.endswith("다.") == True:
                            page_textonly_filtering += line_text + "\n"
                        else:
                            page_textonly_filtering += line_text

                if page_exist_tbl is False:
                    page_textonly_filtering = re.sub(r"(?<![\.\?\!])\n", " ", page_textonly_filtering)
                    whole_text_a_page += page_textonly_filtering
            print(f"{page_num} : {whole_text_a_page}")
            whole_page_extractinfo += re.sub(r"\(cid:[0-9]+\)", "", whole_text_a_page)

        print(whole_page_extractinfo)
        print('Go Next document')
        return whole_page_extractinfo

    def iter_doc_blocks(self, parent):
        """
        yield document element type
        :param parent: word paragraph object
        :return: word element value
        """
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        elif isinstance(parent, _Row):
            parent_elm = parent._tr
        else:
            raise ValueError("something's not right")
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    def get_doc_context(self, doc):
        """
        extraction docx file contents from doc object
        :param doc: doc file object
        :return: extracted doc contents
        """
        doc_contents = ''
        for block in self.iter_doc_blocks(doc):
            if isinstance(block, Paragraph):
                doc_contents += block.text + "\n"
            elif isinstance(block, Table):
                for row in block.rows:
                    row_data = []
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            row_data.append(paragraph.text)
                    doc_contents += "|".join(row_data) + "\n"
        return doc_contents

    def get_ppt_context(self, ppt_prs):
        """
         extraction pptx file contents from ppt object
        :param ppt_prs: ppt object
        :return:  extracted ppt contents
        """
        pptx_contents = ''
        for _, slide in enumerate(ppt_prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        pptx_contents += run.text + '\r\n'
        return pptx_contents

    def extract_file_content(self, file_name):
        """
        extraction contents from various file format
        :param file_name: file name(object name) in object storage
        :return: extract text
        """
        try:
            formed_clear_contents = ''
            if not os.path.exists("tmp/minio_file/"):
                os.makedirs("tmp/minio_file/")
            tmp_local_file_name = "tmp/minio_file/" + file_name
            self._minio_client.fget_object(self._bucket_name, file_name, tmp_local_file_name)
        except Exception as e:
            print(f"Error processing {tmp_local_file_name}: {e}")
            return 0, None, str(e)

        f_extension = pathlib.Path(tmp_local_file_name).suffix
        f_extension = f_extension.lower()

        if f_extension.endswith('.pdf'):
            formed_clear_contents = self.get_context_pdffile_by_plumber(tmp_local_file_name)
            print(formed_clear_contents)

        elif f_extension.endswith('.hwp'):
            hwp_obj = HWPExtractor(tmp_local_file_name)
            hwp_text = hwp_obj.get_text()
            formed_clear_contents = hwp_text

        elif f_extension.endswith('.docx') or f_extension.endswith('.doc'):
            doc = Document(tmp_local_file_name)
            formed_clear_contents = self.get_doc_context(doc)

        elif f_extension.endswith('.txt'):
            loader = TextLoader(tmp_local_file_name)
            docs = loader.load()
            for page in docs:
                formed_clear_contents += page.page_content

        elif f_extension.endswith('.xlsx') or f_extension.endswith('.xls'):
            df = pd.read_excel(tmp_local_file_name)
            df_markdown = df.to_markdown()
            formed_clear_contents = df_markdown

        elif f_extension.endswith('.csv'):
            df = pd.read_csv(tmp_local_file_name)
            df_markdown = df.to_markdown()
            formed_clear_contents = df_markdown

        elif f_extension.endswith('.pptx') or f_extension.endswith('.ppt'):
            try:
                prs = Presentation(tmp_local_file_name)
                formed_clear_contents = self.get_ppt_context(prs)
                print(formed_clear_contents)
            except Exception as e:
                print(e)
        else:
            print("Error: invalid file type")
            print(tmp_local_file_name)
        return 1, formed_clear_contents, "ok"

    def extract_all(self):
        """
        extract contents from all object in target bucket
        :return:
        """
        for item in self._minio_client.list_objects(self._bucket_name, recursive=True):
            if item.is_dir is False:
                print(item.object_name)
                object_name = item.object_name
                path_depth = object_name.split("/")
                if len(path_depth) == 1:
                    print("Warning,...There is no field structure")
                    area = ''
                else:
                    area = path_depth[0]
                if os.path.exists(f"result/{area}") is False:
                    os.makedirs(f"result/{area}")

                _, contents, status = self.extract_file_content(item.object_name)
                if status == "ok":
                    file_fullpath = item.object_name
                    filename = file_fullpath.split("/")[-1]
                    _, tail = os.path.split(filename)
                    tail = tail+".txt"
                    tail = f"result/{area}/"+ tail
                    with open(tail, "w", encoding="utf-8") as fw:
                        print(contents)
                        fw.write(contents)

te = TextExtract(bucket_name="f5e4aeaf4ad899a3a0bf79fea05b7b96820b9103")
te.del_table = True
te.extract_all()
