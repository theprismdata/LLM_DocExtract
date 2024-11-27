import re

import yaml
from minio import Minio
import pdfplumber
from langchain.document_loaders import UnstructuredWordDocumentLoader
from langchain.document_loaders import TextLoader
# from docx import Document
from pptx import Presentation
from HwpParser import HWPExtractor
import pandas as pd
import pathlib
import os

from docx import Document
from docx.document import Document as _Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import _Cell, Table, _Row
from docx.text.paragraph import Paragraph

class TextExtract:
    def __init__(self, bucket_name):
        with open('config/set-dev.yaml') as f:
            config = yaml.load(f, Loader=yaml.FullLoader)
            minio_info = config['minio']
        self.minio_address =  minio_info['address']
        self.accesskey = minio_info['accesskey']
        self.secretkey = minio_info['secretkey']
        self.minio_client = Minio(self.minio_address,
                       access_key=self.accesskey,
                       secret_key=self.secretkey, secure=False)
        self.bucket_name = bucket_name

    def get_pdfpage_info_by_plumber(self, download_file_name)-> list:
        print(f"Source path {download_file_name}")
        try:
            pdfplumb = pdfplumber.open(download_file_name)
            file_extract_contents = ""

            for page_num, plumb_page in enumerate(pdfplumb.pages):
                page_plumb_contents = {}

                table_list = []
                for table_info in pdfplumb.pages[page_num].find_tables():
                    x0 = table_info.bbox[0]
                    y0 = table_info.bbox[1]
                    x1 = table_info.bbox[2]
                    y1 = table_info.bbox[3]
                    table_list.append((x0, y0, x1, y1))
                    table = table_info.extract()
                    df = pd.DataFrame(table[1::], columns=table[0])
                    df.replace('\x00', '', inplace=True)
                    df.replace('Ÿ', '*', inplace=True)
                    page_plumb_contents[int(y0)] = df.to_markdown()

                for content in pdfplumb.pages[page_num].extract_text_lines():
                    x0 = content['x0']
                    y0 = content['top']
                    x1 = content['x1']
                    y1 = content['bottom']
                    if len(table_list) > 0:
                        if table_list[0][0] < x0 and table_list[0][1] < y0 and table_list[0][2] > x1 and table_list[0][3] > y1:
                            pass
                        else:
                            page_plumb_contents[int(y0)] = content['text']
                    else:
                        page_plumb_contents[int(y0)] = content['text']

                if len(page_plumb_contents) > 0:
                    pos_list = list(page_plumb_contents.keys())
                    pos_list = sorted(pos_list)
                    for position in pos_list:
                        line_contents = page_plumb_contents[position]
                        file_extract_contents += line_contents + "\n"

            file_extract_contents = re.sub(r"(?<![\.\?\!])\n", " ", file_extract_contents)
            file_extract_contents = re.sub(r"\(cid:[0-9]+\)", "", file_extract_contents)
            return file_extract_contents
        except Exception as e:
            print(e)

    def iter_block_items(self, parent):
        """
        Generate a reference to each paragraph and table child within *parent*,
        in document order. Each returned value is an instance of either Table or
        Paragraph. *parent* would most commonly be a reference to a main
        Document object, but also works for a _Cell object, which itself can
        contain paragraphs and tables.
        """
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        elif isinstance(parent, _Row):
            parent_elm = parent._tr
        else:
            raise ValueError("something's not right")
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    def extract_file_content(self, file_name):
        try:
            formed_clear_contents = ''
            if not os.path.exists("tmp/minio_file/"):
                os.makedirs("tmp/minio_file/")
            download_file_name = "tmp/minio_file/" + file_name
            self.minio_client.fget_object(self.bucket_name, file_name, download_file_name)

            f_extension = pathlib.Path(download_file_name).suffix
            f_extension = f_extension.lower()

            if f_extension.endswith('.pdf'):
                formed_clear_contents = self.get_pdfpage_info_by_plumber(download_file_name)
                print(formed_clear_contents)

            elif f_extension.endswith('.hwp'):
                hwp_obj = HWPExtractor(download_file_name)
                hwp_text = hwp_obj.get_text()
                formed_clear_contents = hwp_text

            elif f_extension.endswith('.docx') or f_extension.endswith('.doc'):
                doc = Document(download_file_name)
                for block in self.iter_block_items(doc):
                    if isinstance(block, Paragraph):
                        formed_clear_contents += block.text + "\n"
                    elif isinstance(block, Table):
                        for row in block.rows:
                            row_data = []
                            for cell in row.cells:
                                for paragraph in cell.paragraphs:
                                    row_data.append(paragraph.text)
                            formed_clear_contents += "|".join(row_data)+ "\n"

            elif f_extension.endswith('.txt'):
                loader = TextLoader(download_file_name)
                docs = loader.load()
                for page in docs:
                    formed_clear_contents += page.page_content

            elif f_extension.endswith('.xlsx') or f_extension.endswith('.xls'):
                df = pd.read_excel(download_file_name)
                df_markdown = df.to_markdown()
                formed_clear_contents = df_markdown

            elif f_extension.endswith('.csv'):
                df = pd.read_csv(download_file_name)
                df_markdown = df.to_markdown()
                formed_clear_contents = df_markdown

            elif f_extension.endswith('.pptx') or f_extension.endswith('.ppt'):
                try:
                    prs = Presentation(download_file_name)
                    for idx, slide in enumerate(prs.slides):
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    formed_clear_contents += run.text + '\r\n'
                    print(formed_clear_contents)
                except Exception as e:
                    print(e)

            else:
                print("Error: invalid file type")
                print(download_file_name)

        except Exception as e:
            print(f"Error processing {download_file_name}: {e}")
            return 0, None, str(e)
        return 1, formed_clear_contents, "ok"

    def extract_all(self):
        for item in self.minio_client.list_objects(self.bucket_name,recursive=True):
            if item.is_dir is False:
                print(item.object_name)
                object_name = item.object_name
                path_depth = object_name.split("/")
                if len(path_depth) == 1:
                    print("Warning,...There is no field structure")
                    area = ''
                else:
                    area = path_depth[0]
                if os.path.exists(f"result/{area}") == False:
                    os.makedirs(f"result/{area}")

                _, contents, status = self.extract_file_content(item.object_name)
                if status == "ok":
                    file_fullpath = item.object_name
                    filename = file_fullpath.split("/")[-1]
                    head, tail = os.path.split(filename)
                    tail = tail+".txt"
                    tail = f"result/{area}/"+ tail
                    with open(tail, "w") as fw:
                        print(contents)
                        fw.write(contents)

te = TextExtract(bucket_name="guest003")
te.extract_all()
